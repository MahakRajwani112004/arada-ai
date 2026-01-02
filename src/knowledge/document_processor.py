"""Enhanced document processor for knowledge base ingestion.

Supports comprehensive extraction from various document types:
- PDF: Text, tables, images (with OCR for scanned docs)
- DOCX: Text, tables, images
- PPTX: Slides, speaker notes, images
- Excel: Sheets as tables
- Images: OCR text extraction
- Text/Markdown: Direct text

Uses pymupdf4llm for PDF-to-Markdown conversion optimized for LLM/RAG workflows.
"""
import io
import os
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from typing import Any, Dict, List, Optional, Tuple

from src.config.logging import get_logger
from src.config.settings import get_settings

logger = get_logger(__name__)


class DocumentType(str, Enum):
    """Supported document types."""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    CSV = "csv"
    TXT = "txt"
    MD = "md"
    HTML = "html"
    IMAGE = "image"  # jpg, jpeg, png, gif, bmp, tiff


# File extension to document type mapping
EXTENSION_MAP: Dict[str, DocumentType] = {
    "pdf": DocumentType.PDF,
    "docx": DocumentType.DOCX,
    "doc": DocumentType.DOCX,  # Will attempt conversion
    "pptx": DocumentType.PPTX,
    "ppt": DocumentType.PPTX,
    "xlsx": DocumentType.XLSX,
    "xls": DocumentType.XLSX,
    "csv": DocumentType.CSV,
    "txt": DocumentType.TXT,
    "md": DocumentType.MD,
    "markdown": DocumentType.MD,
    "html": DocumentType.HTML,
    "htm": DocumentType.HTML,
    "jpg": DocumentType.IMAGE,
    "jpeg": DocumentType.IMAGE,
    "png": DocumentType.IMAGE,
    "gif": DocumentType.IMAGE,
    "bmp": DocumentType.IMAGE,
    "tiff": DocumentType.IMAGE,
    "tif": DocumentType.IMAGE,
    "webp": DocumentType.IMAGE,
}

# MIME type mapping
MIME_TYPE_MAP: Dict[str, str] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "csv": "text/csv",
    "txt": "text/plain",
    "md": "text/markdown",
    "html": "text/html",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "png": "image/png",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tiff": "image/tiff",
    "webp": "image/webp",
}


@dataclass
class ExtractedImage:
    """Represents an extracted image from a document."""
    data: bytes
    format: str  # jpg, png, etc.
    page_number: Optional[int] = None
    description: Optional[str] = None  # AI-generated description
    width: Optional[int] = None
    height: Optional[int] = None


@dataclass
class ExtractedTable:
    """Represents an extracted table from a document."""
    markdown: str  # Table in markdown format
    page_number: Optional[int] = None
    sheet_name: Optional[str] = None  # For Excel files
    row_count: int = 0
    col_count: int = 0


@dataclass
class ExtractionResult:
    """Result of document extraction."""
    text: str  # Main extracted text (may include table markdown)
    tables: List[ExtractedTable] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    page_count: int = 0
    word_count: int = 0
    has_ocr_content: bool = False
    error: Optional[str] = None


class BaseDocumentProcessor(ABC):
    """Abstract base class for document processors."""

    @abstractmethod
    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from a document."""
        pass

    @abstractmethod
    def supports(self, doc_type: DocumentType) -> bool:
        """Check if this processor supports the given document type."""
        pass


class PDFProcessor(BaseDocumentProcessor):
    """PDF processor using pymupdf4llm for LLM-optimized extraction."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.PDF

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from PDF using pymupdf4llm."""
        try:
            import fitz  # PyMuPDF
            import pymupdf4llm
        except ImportError as e:
            return ExtractionResult(
                text="",
                error=f"PDF processing libraries not installed: {e}",
            )

        try:
            # Open document
            doc = fitz.open(file_path)
            page_count = len(doc)

            # Extract metadata
            metadata = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
                "creator": doc.metadata.get("creator", ""),
                "producer": doc.metadata.get("producer", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
                "modification_date": doc.metadata.get("modDate", ""),
            }

            # Use pymupdf4llm for markdown extraction (includes tables)
            # This is optimized for LLM/RAG workflows
            md_text = pymupdf4llm.to_markdown(
                doc,
                page_chunks=False,  # Return as single text
                write_images=False,  # We'll extract images separately
                show_progress=False,
            )

            # Check if document might be scanned (very little text)
            has_ocr_content = False
            if len(md_text.strip()) < 100 and page_count > 0:
                # Try OCR extraction
                ocr_text = await self._extract_with_ocr(doc)
                if ocr_text:
                    md_text = ocr_text
                    has_ocr_content = True

            # Extract tables separately for structured access
            tables = await self._extract_tables(doc)

            # Extract images
            images = await self._extract_images(doc)

            doc.close()

            # Calculate word count
            word_count = len(md_text.split())

            return ExtractionResult(
                text=md_text,
                tables=tables,
                images=images,
                metadata=metadata,
                page_count=page_count,
                word_count=word_count,
                has_ocr_content=has_ocr_content,
            )

        except Exception as e:
            logger.error("pdf_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))

    async def _extract_with_ocr(self, doc) -> str:
        """Extract text using OCR for scanned documents."""
        try:
            import fitz  # PyMuPDF - needed for Matrix
            import pytesseract
            from PIL import Image
        except ImportError:
            logger.warning("OCR libraries not available")
            return ""

        text_parts = []
        for page_num, page in enumerate(doc):
            # Render page to image
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
            pix = page.get_pixmap(matrix=mat)

            # Convert to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # OCR
            try:
                page_text = pytesseract.image_to_string(img, lang='eng')
                if page_text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
            except Exception as e:
                logger.warning("ocr_page_failed", page=page_num, error=str(e))

        return "\n\n".join(text_parts)

    async def _extract_tables(self, doc) -> List[ExtractedTable]:
        """Extract tables from PDF."""
        tables = []
        for page_num, page in enumerate(doc):
            try:
                # Find tables on the page
                page_tables = page.find_tables()
                for table in page_tables:
                    md = table.to_markdown()
                    if md.strip():
                        tables.append(ExtractedTable(
                            markdown=md,
                            page_number=page_num + 1,
                            row_count=len(table.extract()),
                            col_count=len(table.extract()[0]) if table.extract() else 0,
                        ))
            except Exception as e:
                logger.warning("table_extraction_failed", page=page_num, error=str(e))
        return tables

    async def _extract_images(self, doc) -> List[ExtractedImage]:
        """Extract images from PDF."""
        images = []
        for page_num, page in enumerate(doc):
            try:
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        if base_image:
                            images.append(ExtractedImage(
                                data=base_image["image"],
                                format=base_image["ext"],
                                page_number=page_num + 1,
                                width=base_image.get("width"),
                                height=base_image.get("height"),
                            ))
                    except Exception as e:
                        logger.warning(
                            "image_extraction_failed",
                            page=page_num,
                            index=img_index,
                            error=str(e),
                        )
            except Exception as e:
                logger.warning("page_image_list_failed", page=page_num, error=str(e))
        return images


class DOCXProcessor(BaseDocumentProcessor):
    """DOCX processor with table and image extraction."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.DOCX

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from DOCX."""
        try:
            from docx import Document
            from docx.table import Table
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
        except ImportError as e:
            return ExtractionResult(
                text="",
                error=f"python-docx not installed: {e}",
            )

        try:
            doc = Document(file_path)
            text_parts = []
            tables = []
            images = []

            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
            }

            # Process paragraphs and tables in order
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    # Paragraph
                    for para in doc.paragraphs:
                        if para._element == element and para.text.strip():
                            text_parts.append(para.text)
                            break
                elif element.tag.endswith('tbl'):
                    # Table
                    for table in doc.tables:
                        if table._element == element:
                            md_table = self._table_to_markdown(table)
                            tables.append(ExtractedTable(
                                markdown=md_table,
                                row_count=len(table.rows),
                                col_count=len(table.columns),
                            ))
                            text_parts.append(md_table)
                            break

            # Extract images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        # Determine format from content type
                        content_type = rel.target_part.content_type
                        ext = content_type.split("/")[-1] if "/" in content_type else "png"
                        images.append(ExtractedImage(
                            data=image_data,
                            format=ext,
                        ))
                    except Exception as e:
                        logger.warning("docx_image_extraction_failed", error=str(e))

            text = "\n\n".join(text_parts)
            word_count = len(text.split())

            return ExtractionResult(
                text=text,
                tables=tables,
                images=images,
                metadata=metadata,
                page_count=1,  # DOCX doesn't have fixed pages
                word_count=word_count,
            )

        except Exception as e:
            logger.error("docx_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))

    def _table_to_markdown(self, table) -> str:
        """Convert DOCX table to markdown format."""
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                # Add header separator
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(rows)


class PPTXProcessor(BaseDocumentProcessor):
    """PowerPoint processor with slide and speaker notes extraction."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.PPTX

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from PowerPoint."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            return ExtractionResult(
                text="",
                error=f"python-pptx not installed: {e}",
            )

        try:
            prs = Presentation(file_path)
            text_parts = []
            images = []

            metadata = {
                "slide_count": len(prs.slides),
            }

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"## Slide {slide_num}")

                # Extract text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_text.append(text)

                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            images.append(ExtractedImage(
                                data=image.blob,
                                format=image.ext,
                                page_number=slide_num,
                            ))
                        except Exception as e:
                            logger.warning(
                                "pptx_image_extraction_failed",
                                slide=slide_num,
                                error=str(e),
                            )

                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_text.append(f"\n**Speaker Notes:** {notes_frame.text}")

                text_parts.append("\n".join(slide_text))

            text = "\n\n".join(text_parts)
            word_count = len(text.split())

            return ExtractionResult(
                text=text,
                images=images,
                metadata=metadata,
                page_count=len(prs.slides),
                word_count=word_count,
            )

        except Exception as e:
            logger.error("pptx_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))


class ExcelProcessor(BaseDocumentProcessor):
    """Excel processor with sheet-to-table conversion."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type in [DocumentType.XLSX, DocumentType.CSV]

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from Excel/CSV."""
        try:
            import pandas as pd
        except ImportError as e:
            return ExtractionResult(
                text="",
                error=f"pandas not installed: {e}",
            )

        try:
            doc_type = get_document_type(filename)
            tables = []
            text_parts = []

            if doc_type == DocumentType.CSV:
                # Single CSV file
                df = pd.read_csv(file_path)
                md = df.to_markdown(index=False)
                tables.append(ExtractedTable(
                    markdown=md,
                    sheet_name="Sheet1",
                    row_count=len(df),
                    col_count=len(df.columns),
                ))
                text_parts.append(f"## Sheet: Sheet1\n\n{md}")
            else:
                # Excel with multiple sheets
                try:
                    import openpyxl
                    xls = pd.ExcelFile(file_path, engine='openpyxl')
                except ImportError:
                    xls = pd.ExcelFile(file_path)

                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    if not df.empty:
                        md = df.to_markdown(index=False)
                        tables.append(ExtractedTable(
                            markdown=md,
                            sheet_name=sheet_name,
                            row_count=len(df),
                            col_count=len(df.columns),
                        ))
                        text_parts.append(f"## Sheet: {sheet_name}\n\n{md}")

            text = "\n\n".join(text_parts)
            word_count = len(text.split())

            return ExtractionResult(
                text=text,
                tables=tables,
                metadata={"sheet_count": len(tables)},
                page_count=len(tables),
                word_count=word_count,
            )

        except Exception as e:
            logger.error("excel_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))


class ImageProcessor(BaseDocumentProcessor):
    """Image processor using OCR."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.IMAGE

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract text from image using OCR."""
        try:
            import pytesseract
            from PIL import Image
        except ImportError as e:
            return ExtractionResult(
                text="",
                error=f"OCR libraries not installed: {e}",
            )

        try:
            img = Image.open(file_path)

            # Get image metadata
            metadata = {
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "mode": img.mode,
            }

            # Perform OCR
            text = pytesseract.image_to_string(img, lang='eng')

            # Store the image itself
            with open(file_path, 'rb') as f:
                image_data = f.read()

            ext = filename.rsplit(".", 1)[-1].lower()
            images = [ExtractedImage(
                data=image_data,
                format=ext,
                width=img.width,
                height=img.height,
            )]

            word_count = len(text.split())

            return ExtractionResult(
                text=text,
                images=images,
                metadata=metadata,
                page_count=1,
                word_count=word_count,
                has_ocr_content=True,
            )

        except Exception as e:
            logger.error("image_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))


class TextProcessor(BaseDocumentProcessor):
    """Plain text and markdown processor."""

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type in [DocumentType.TXT, DocumentType.MD, DocumentType.HTML]

    async def extract(self, file_path: str, filename: str) -> ExtractionResult:
        """Extract content from text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()

            doc_type = get_document_type(filename)

            # For HTML, try to extract just the text
            if doc_type == DocumentType.HTML:
                text = self._strip_html(text)

            word_count = len(text.split())

            return ExtractionResult(
                text=text,
                metadata={"encoding": "utf-8"},
                page_count=1,
                word_count=word_count,
            )

        except Exception as e:
            logger.error("text_extraction_failed", filename=filename, error=str(e))
            return ExtractionResult(text="", error=str(e))

    def _strip_html(self, html: str) -> str:
        """Strip HTML tags and extract text."""
        try:
            from html.parser import HTMLParser

            class HTMLTextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                    self.skip_data = False

                def handle_starttag(self, tag, attrs):
                    if tag in ['script', 'style', 'head']:
                        self.skip_data = True

                def handle_endtag(self, tag):
                    if tag in ['script', 'style', 'head']:
                        self.skip_data = False

                def handle_data(self, data):
                    if not self.skip_data:
                        text = data.strip()
                        if text:
                            self.text_parts.append(text)

            parser = HTMLTextExtractor()
            parser.feed(html)
            return "\n".join(parser.text_parts)
        except Exception:
            # Fallback: simple regex
            return re.sub(r'<[^>]+>', '', html)


# Registry of processors
_PROCESSORS: List[BaseDocumentProcessor] = [
    PDFProcessor(),
    DOCXProcessor(),
    PPTXProcessor(),
    ExcelProcessor(),
    ImageProcessor(),
    TextProcessor(),
]


def get_document_type(filename: str) -> Optional[DocumentType]:
    """Get document type from filename extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return EXTENSION_MAP.get(ext)


def get_supported_types() -> List[str]:
    """Get list of supported file extensions."""
    return list(EXTENSION_MAP.keys())


def get_mime_type(filename: str) -> str:
    """Get MIME type for a filename."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return MIME_TYPE_MAP.get(ext, "application/octet-stream")


async def process_document(file_path: str, filename: str) -> ExtractionResult:
    """
    Process a document and extract its content.

    Args:
        file_path: Path to the document file
        filename: Original filename (used for type detection)

    Returns:
        ExtractionResult with extracted text, tables, images, and metadata
    """
    if not os.path.exists(file_path):
        return ExtractionResult(
            text="",
            error=f"File not found: {file_path}",
        )

    doc_type = get_document_type(filename)
    if not doc_type:
        return ExtractionResult(
            text="",
            error=f"Unsupported file type. Supported: {', '.join(get_supported_types())}",
        )

    # Find appropriate processor
    for processor in _PROCESSORS:
        if processor.supports(doc_type):
            logger.info(
                "document_processing_started",
                filename=filename,
                doc_type=doc_type.value,
                processor=processor.__class__.__name__,
            )
            result = await processor.extract(file_path, filename)
            logger.info(
                "document_processing_completed",
                filename=filename,
                doc_type=doc_type.value,
                word_count=result.word_count,
                tables=len(result.tables),
                images=len(result.images),
                has_ocr=result.has_ocr_content,
                error=result.error,
            )
            return result

    return ExtractionResult(
        text="",
        error=f"No processor found for document type: {doc_type}",
    )


async def describe_images_with_llm(
    images: List[ExtractedImage],
    max_images: int = 5,
) -> List[ExtractedImage]:
    """
    Generate descriptions for images using OpenAI Vision API.

    Args:
        images: List of extracted images
        max_images: Maximum number of images to describe (to control costs)

    Returns:
        Images with descriptions populated
    """
    settings = get_settings()
    if not settings.openai_api_key:
        logger.warning("OpenAI API key not set, skipping image description")
        return images

    try:
        from openai import AsyncOpenAI
        import base64
    except ImportError:
        return images

    client = AsyncOpenAI(api_key=settings.openai_api_key)

    for i, image in enumerate(images[:max_images]):
        try:
            # Convert image to base64
            b64_image = base64.b64encode(image.data).decode('utf-8')
            mime_type = f"image/{image.format}"

            response = await client.chat.completions.create(
                model="gpt-4o-mini",  # Use mini for cost efficiency
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this image in 1-2 sentences for use in a document search index. Focus on the key information or content shown.",
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{b64_image}",
                                    "detail": "low",  # Use low detail for faster/cheaper processing
                                },
                            },
                        ],
                    }
                ],
                max_tokens=100,
            )

            image.description = response.choices[0].message.content
            logger.debug("image_described", index=i, description=image.description[:50])

        except Exception as e:
            logger.warning("image_description_failed", index=i, error=str(e))

    return images
